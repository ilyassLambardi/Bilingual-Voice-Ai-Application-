"""
Generate defense presentation as editable PowerPoint (.pptx).
Run: python scripts/generate_pptx.py
Output: Bilingual_AI_Defense.pptx

Each slide has ONE unique angle — zero repetition between slides.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Theme colors ──────────────────────────────────────────────────────
NAVY = RGBColor(0x1B, 0x1F, 0x3B)
ACCENT = RGBColor(0x2D, 0x3A, 0x8C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF4, 0xF4, 0xF8)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x6B, 0x6B, 0x7B)
HIGHLIGHT = RGBColor(0x3D, 0x5A, 0xFE)
GREEN = RGBColor(0x1B, 0x8C, 0x4E)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_shape(slide, text, left, top, width, height, font_size=32, bold=True, color=NAVY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    return tf


def add_body_text(slide, text, left, top, width, height, font_size=16, color=DARK_TEXT, bold=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    return tf


def add_bullet_list(slide, items, left, top, width, height, font_size=14, color=DARK_TEXT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"\u2022 {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return tf


def add_table(slide, rows, cols, data, left, top, width, height):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    for j, val in enumerate(data[0]):
        cell = table.cell(0, j)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    for i in range(1, rows):
        for j in range(cols):
            cell = table.cell(i, j)
            cell.text = data[i][j]
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_TEXT
            p.alignment = PP_ALIGN.CENTER
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG

    return table


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # ════════════════════════════════════════════════════════════════════
    # SLIDE 1 — Title
    # UNIQUE ANGLE: Who, what, where — first impression only
    # ════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_title_shape(slide,
        "Design and Implementation of An On-Device\nBilingual AI Voice Assistant System\nfor English and German Based on Deep Learning",
        Inches(1), Inches(1.2), Inches(11), Inches(2.5),
        font_size=36, color=NAVY)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.8), Inches(4), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

    info_lines = [
        "Ilyass Lambardi (202239060075)",
        "Supervisor: Jianbo Wang",
        "School of Computer Science and Software Engineering",
        "Southwest Petroleum University — May 2026",
    ]
    add_bullet_list(slide, info_lines, Inches(1), Inches(4.2), Inches(8), Inches(2.5),
                    font_size=16, color=GRAY)

    # ════════════════════════════════════════════════════════════════════
    # SLIDE 2 — WHY: The Gap in Existing Systems
    # UNIQUE ANGLE: Problem framing only. No solution details.
    # ════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_title_shape(slide, "Why This Project?",
        Inches(0.5), Inches(0.3), Inches(8), Inches(0.8), font_size=32)

    add_body_text(slide, "What bilingual users face today:",
                  Inches(0.5), Inches(1.3), Inches(6), Inches(0.4),
                  font_size=16, color=GRAY, bold=True)

    # Problem scenarios (storytelling approach — not bullet-point lists)
    scenarios = [
        '"Hey Siri, switch to German" — you must TELL the system your language every time',
        '"Ich meine, that\'s just how it works" — mixed speech? System crashes or ignores it',
        "3-5 seconds of silence after you speak — waiting for cloud round-trip + full generation",
        "All your voice data leaves your device — no offline option exists",
    ]
    add_bullet_list(slide, scenarios, Inches(0.5), Inches(2.0), Inches(12), Inches(3),
                    font_size=14, color=DARK_TEXT)

    add_body_text(slide, "No existing open-source system solves all four problems simultaneously.",
                  Inches(0.5), Inches(5.2), Inches(10), Inches(0.5),
                  font_size=15, color=ACCENT, bold=True)

    # ════════════════════════════════════════════════════════════════════
    # SLIDE 3 — WHAT: High-Level Design (the "napkin sketch")
    # UNIQUE ANGLE: Bird's eye view only. Model names, no internals.
    # ════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_title_shape(slide, "What I Built",
        Inches(0.5), Inches(0.3), Inches(8), Inches(0.8), font_size=32)

    # Pipeline boxes — simplified, no model names (those go in slide 5)
    stages = ["You Speak", "Detect\nSpeech", "Transcribe", "Think &\nRespond", "Synthesize\nVoice", "You Hear"]
    box_w = Inches(1.6)
    box_h = Inches(0.9)
    start_x = Inches(0.5)
    y = Inches(1.6)
    gap = Inches(0.35)

    for i, s in enumerate(stages):
        x = start_x + i * (box_w + gap)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT if i in (1, 2, 3, 4) else LIGHT_BG
        shape.line.color.rgb = ACCENT
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = s
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE if i in (1, 2, 3, 4) else NAVY
        p.alignment = PP_ALIGN.CENTER

        if i < len(stages) - 1:
            arrow_x = x + box_w
            arrow_y = y + box_h / 2
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, arrow_y - Pt(8), gap, Pt(16))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT
            arrow.line.fill.background()

    # Three unique selling points (not repeated elsewhere)
    add_body_text(slide, "Three design principles that make this different:",
                  Inches(0.5), Inches(3.0), Inches(8), Inches(0.4),
                  font_size=14, color=GRAY, bold=True)

    principles = [
        "Stream, don't batch — user hears the first sentence while the AI is still thinking",
        "Detect, don't ask — language is inferred from what you say, never toggled manually",
        "Two roads, one interface — cloud for quality, local for privacy, same user experience",
    ]
    add_bullet_list(slide, principles, Inches(0.5), Inches(3.5), Inches(12), Inches(2.5),
                    font_size=14, color=DARK_TEXT)

    # ════════════════════════════════════════════════════════════════════
    # SLIDE 4 — HOW (1): The Streaming Trick
    # UNIQUE ANGLE: Only the latency innovation. No module names.
    # ════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_title_shape(slide, "How: Sentence-Level Streaming",
        Inches(0.5), Inches(0.3), Inches(10), Inches(0.8), font_size=30)

    # Visual timeline comparison
    add_body_text(slide, "Before (traditional):",
                  Inches(0.5), Inches(1.4), Inches(4), Inches(0.4),
                  font_size=14, color=GRAY, bold=True)
    # Gray bar representing wait
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.9), Inches(10), Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    bar.line.fill.background()
    tf = bar.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "Generate ALL tokens ─────────────────── Synthesize ALL ─── Play"
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

    add_body_text(slide, "User waits 4-8 seconds before hearing anything",
                  Inches(0.5), Inches(2.5), Inches(8), Inches(0.3),
                  font_size=11, color=GRAY)

    add_body_text(slide, "After (our system):",
                  Inches(0.5), Inches(3.1), Inches(4), Inches(0.4),
                  font_size=14, color=ACCENT, bold=True)

    # Overlapping colored bars
    bar1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.6), Inches(3.5), Inches(0.4))
    bar1.fill.solid()
    bar1.fill.fore_color.rgb = ACCENT
    bar1.line.fill.background()
    tf = bar1.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "Sentence 1 → TTS → PLAY"
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    bar2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(4.1), Inches(3.5), Inches(0.4))
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = HIGHLIGHT
    bar2.line.fill.background()
    tf = bar2.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "Sentence 2 → TTS → queue"
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    bar3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(4.6), Inches(3.5), Inches(0.4))
    bar3.fill.solid()
    bar3.fill.fore_color.rgb = RGBColor(0x5C, 0x6B, 0xC0)
    bar3.line.fill.background()
    tf = bar3.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "Sentence 3 → TTS → queue"
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    add_body_text(slide, "User hears audio after ~1.5s (first sentence only). Rest overlaps.",
                  Inches(0.5), Inches(5.3), Inches(10), Inches(0.4),
                  font_size=14, color=GREEN, bold=True)

    add_body_text(slide, "Result: ~40% perceived latency reduction. Feels instant in conversation.",
                  Inches(0.5), Inches(5.8), Inches(10), Inches(0.4),
                  font_size=13, color=DARK_TEXT)

    # ════════════════════════════════════════════════════════════════════
    # SLIDE 5 — HOW (2): What Happens When You Interrupt
    # UNIQUE ANGLE: Only interruption + backchannel logic. Fresh content.
    # ════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_title_shape(slide, "How: Natural Turn-Taking",
        Inches(0.5), Inches(0.3), Inches(10), Inches(0.8), font_size=30)

    add_body_text(slide, "Problem: How does the system know when to stop vs. when to keep talking?",
                  Inches(0.5), Inches(1.2), Inches(11), Inches(0.4),
                  font_size=14, color=GRAY)

    # Two-column layout
    add_body_text(slide, "Interrupt (stop AI)", Inches(0.5), Inches(2.0), Inches(5), Inches(0.4),
                  font_size=16, color=ACCENT, bold=True)
    interrupt_info = [
        "User speaks for 4+ frames (~128ms) while AI is talking",
        "VAD detects sustained speech → sets interrupt flag",
        "Generation task cancelled, audio stops immediately",
        "Pipeline resets to Listening state for new input",
    ]
    add_bullet_list(slide, interrupt_info, Inches(0.5), Inches(2.5), Inches(5.5), Inches(2.5),
                    font_size=13, color=DARK_TEXT)

    add_body_text(slide, "Backchannel (ignore)", Inches(7), Inches(2.0), Inches(5), Inches(0.4),
                  font_size=16, color=GREEN, bold=True)
    backchannel_info = [
        "User makes short sound < 384ms ('mhm', 'yeah', 'ok')",
        "VAD detects brief burst → classifies as backchannel",
        "AI continues speaking without interruption",
        "Feels natural — like talking to a real person",
    ]
    add_bullet_list(slide, backchannel_info, Inches(7), Inches(2.5), Inches(5.5), Inches(2.5),
                    font_size=13, color=DARK_TEXT)

    add_body_text(slide, "This is what makes the system feel conversational, not robotic.",
                  Inches(0.5), Inches(5.5), Inches(10), Inches(0.4),
                  font_size=14, color=NAVY, bold=True)

    # ════════════════════════════════════════════════════════════════════
    # SLIDE 6 — HOW (3): Language Intelligence (the "brain")
    # UNIQUE ANGLE: Only the bilingual detection logic. No model specs.
    # ════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_title_shape(slide, "How: Bilingual Without a Button",
        Inches(0.5), Inches(0.3), Inches(10), Inches(0.8), font_size=30)

    add_body_text(slide, "The Challenge:", Inches(0.5), Inches(1.2), Inches(8), Inches(0.4),
                  font_size=14, color=GRAY)
    add_body_text(slide, "ASR outputs a language tag, but it's often WRONG for short or mixed utterances.",
                  Inches(0.5), Inches(1.6), Inches(10), Inches(0.4),
                  font_size=14, color=DARK_TEXT)

    add_body_text(slide, "Our Solution: Intent-Based Override", Inches(0.5), Inches(2.3), Inches(8), Inches(0.4),
                  font_size=16, color=ACCENT, bold=True)

    solution_points = [
        "Analyze sentence structure (word order, articles, verb position) to detect true language",
        "Override ASR tag when structural signals contradict it",
        "Track language history (last 5 utterances) to detect deliberate switches vs. noise",
        'Special triggers: "What does X mean?" activates Teacher Mode regardless of language',
    ]
    add_bullet_list(slide, solution_points, Inches(0.5), Inches(2.9), Inches(12), Inches(2),
                    font_size=13, color=DARK_TEXT)

    # Example conversation
    add_body_text(slide, "Live Example:", Inches(0.5), Inches(4.8), Inches(3), Inches(0.4),
                  font_size=14, color=ACCENT, bold=True)
    examples = [
        'You: "Hey, what does Fernweh mean?" → AI detects EN + teacher trigger',
        'AI: "Fernweh means a longing for faraway places..." (responds in English with German examples)',
        'You: "Erzähl mir mehr darüber" → AI detects full DE → switches entirely to German',
    ]
    add_bullet_list(slide, examples, Inches(0.5), Inches(5.3), Inches(12), Inches(1.8),
                    font_size=12, color=DARK_TEXT)

    # ════════════════════════════════════════════════════════════════════
    # SLIDE 7 — HOW (4): Tech Stack (the "what's under the hood" slide)
    # UNIQUE ANGLE: Only the model/tech choices. No behavior description.
    # ════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_title_shape(slide, "Under the Hood: Technology Choices",
        Inches(0.5), Inches(0.3), Inches(10), Inches(0.8), font_size=30)

    data = [
        ["Layer", "Cloud Mode", "Local Mode", "Why This Choice"],
        ["Server", "FastAPI + WebSocket", "Same", "Async, binary streaming support"],
        ["VAD", "Silero v5 (2M params)", "Same", "CPU-only, 32ms frame latency"],
        ["ASR", "Whisper Large v3 (1.5B)", "faster-whisper small", "Best accuracy vs. VRAM trade-off"],
        ["LLM", "Llama 3.3 70B (Groq)", "Qwen 2.5 1.5B", "70B quality or 4GB-VRAM offline"],
        ["TTS", "Edge Neural TTS", "Silero v3 (EN+DE)", "Natural voice vs. zero-latency"],
        ["Memory", "SQLite + embeddings", "Same", "Zero-config, single file, portable"],
    ]
    add_table(slide, 7, 4, data, Inches(0.3), Inches(1.3), Inches(12.5), Inches(4))

    add_body_text(slide, "Key trade-off: Cloud = best quality + multi-user | Local = private + offline + single-user",
                  Inches(0.5), Inches(5.8), Inches(11), Inches(0.4),
                  font_size=13, color=GRAY)

    # ════════════════════════════════════════════════════════════════════
    # SLIDE 8 — RESULTS: Numbers That Matter
    # UNIQUE ANGLE: Only measurable outcomes. No design explanation.
    # ════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_title_shape(slide, "Results",
        Inches(0.5), Inches(0.3), Inches(8), Inches(0.8), font_size=32)

    # Three big metric blocks
    metrics = [
        ("1.5-2.5s", "End-to-end latency\n(vs. 4-8s traditional)"),
        ("~128ms", "Interrupt response\n(vs. not supported)"),
        (">95%", "Language detection\n(EN and DE combined)"),
    ]
    for i, (big, sub) in enumerate(metrics):
        x = Inches(0.5) + i * Inches(4.2)
        add_title_shape(slide, big, x, Inches(1.4), Inches(3.8), Inches(1),
                        font_size=44, color=ACCENT, bold=True)
        add_body_text(slide, sub, x, Inches(2.6), Inches(3.8), Inches(0.8),
                      font_size=13, color=GRAY)

    # What works well
    add_body_text(slide, "What the system handles correctly:", Inches(0.5), Inches(3.8), Inches(6), Inches(0.4),
                  font_size=14, color=DARK_TEXT, bold=True)
    works = [
        "Pure EN or DE recognized with >98%/>95% accuracy",
        "Mid-conversation code-switching without user action",
        "Backchannel sounds ignored; real interrupts caught in ~128ms",
        "Cross-session memory recall (remembers past conversation topics)",
        "Hallucination filter rejects >90% of noise-only transcripts",
    ]
    add_bullet_list(slide, works, Inches(0.5), Inches(4.3), Inches(12), Inches(2.8),
                    font_size=13, color=DARK_TEXT)

    # ════════════════════════════════════════════════════════════════════
    # SLIDE 9 — HONEST: What Doesn't Work + What's Next
    # UNIQUE ANGLE: Candid limitations + clear next steps
    # ════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_title_shape(slide, "Limitations & Next Steps",
        Inches(0.5), Inches(0.3), Inches(8), Inches(0.8), font_size=30)

    add_body_text(slide, "What I'd do differently / next:", Inches(0.5), Inches(1.2), Inches(6), Inches(0.4),
                  font_size=14, color=GRAY)

    # Combined — each limitation paired with its future fix
    pairs = [
        "Only EN + DE → Next: add Chinese, French, Arabic (modular TTS plugins)",
        "Cloud needed for best quality → Next: better local LLMs (Phi-3, Gemma 2)",
        "No voice personalization → Next: XTTSv2 voice cloning from user samples",
        "Desktop only → Next: mobile via ONNX runtime optimization",
        "Simple memory retrieval → Next: RAG with knowledge graph for factual grounding",
    ]
    add_bullet_list(slide, pairs, Inches(0.5), Inches(1.8), Inches(12), Inches(4),
                    font_size=14, color=DARK_TEXT)

    # ════════════════════════════════════════════════════════════════════
    # SLIDE 10 — Demo (last slide before questions)
    # UNIQUE ANGLE: Video + what to watch for
    # ════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_title_shape(slide, "Live Demo",
        Inches(0.5), Inches(0.3), Inches(8), Inches(0.8), font_size=32)

    # Video placeholder
    placeholder = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(1.3), Inches(9), Inches(5))
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = LIGHT_BG
    placeholder.line.color.rgb = ACCENT
    tf = placeholder.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "[1-Minute Screen Recording]"
    p.font.size = Pt(24)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

    add_body_text(slide, "Watch for: language switch \u2022 interrupt \u2022 streaming speed \u2022 orb reacting",
                  Inches(2), Inches(6.5), Inches(9), Inches(0.4),
                  font_size=13, color=ACCENT, bold=True)

    # ── Save ──────────────────────────────────────────────────────────
    output_path = "Bilingual_AI_Defense_v2.pptx"
    prs.save(output_path)
    print(f"Presentation saved: {output_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_presentation()
